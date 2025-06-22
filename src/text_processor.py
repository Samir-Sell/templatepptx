from parent_processor import ParentProcessor
from pptx.presentation import Presentation as PowerPoint
from  pptx.shapes.autoshape import Shape

class TextProcessor(ParentProcessor):

    def __init__(self, shape: Shape, context: dict, slide_number: int, special_character: str):
        super().__init__(shape, context, slide_number, special_character)

    def replace_text(self):

        '''
        Description: Process text and replace every value
        
        @input shape: A container from PPTX Python called shape which contains paragraphs and text
        @input context: A dictionary containing all of the data that is fed into the template. It contains the data 
        and the magic keywords.
        @input slide_number: The slide number index
        '''
    
        if self._shape.has_text_frame:
            for key in self._context:
                if(self._shape.text.find(str(key)))!=-1:
                    text_frame = self._shape.text_frame
                    for p in text_frame.paragraphs:
                        self._replace_runs(p)

    def _formatting_check(self, run, last_run) -> bool:

        """
        description: Checks for the formatting of the different text runs and compares them

        @input run: Run object that contains formatting and text for a small segment of text
        @input last_run: Run object that contains formatting and text for a small segment of text from the last run
        
        @output: Boolean containing True if current run matches last run
        """
        
        # Compare two runs formatting
        if last_run == None:
            return False
        run_format = {'size':run.font.size,
                    'bold':run.font.bold,
                    'underline':run.font.underline,
                    'italic':run.font.italic,
                    'name' :run.font.name}
        last_run_format = {'size':last_run.font.size,
                    'bold':last_run.font.bold,
                    'underline':last_run.font.underline,
                    'italic':last_run.font.italic,
                    'name' :last_run.font.name}
        if run_format == last_run_format:
            return True
        else:
            return False
    
    def _replace_runs(self, p):
        
        """
        Description: Function to replace runs for each paragraph object.
        
        @input p: Paragraph object from pptx that contains multiple runs
        @input slide_number: integer containing the slide number
        @input context: A dictionary containing all of the data that is fed into the template. It contains the data 
        and the magic keywords.
        """

        last_run = None
        # Loop through every run in a paragraph obj
        for run in p.runs:
            # Check if the current run and last have the same formatting. If they do, combine the runs
            if self._formatting_check(run, last_run):
                combined_text = last_run.text + run.text
                for key in self._context:
                    to_find = self._add_special(key)
                    combined_text = combined_text.replace(str(to_find), str(self._context[key]))
                run.text = combined_text
                # Remove the current run after it has been combined with the last run
                run_to_remove = last_run._r
                run_to_remove.getparent().remove(run_to_remove)
            # If formatting does not align
            else:
                # Replace values in the run
                text = run.text
                for key in self._context:
                    to_find = self._add_special(str(key))
                    text = text.replace(str(to_find), str(self._context[key]))
                run.text = text
            last_run = run

    def _add_special(self, key) -> str:
        
        """
        Description: Adds a 'magic' character to either side of the dictionary key being handed to the function

        @input key: A key from the context dictionary that acts as a magic word in the PPTX template
        @input SPECIAL_CHARACTER: A character that is added to the front and the end of the key
        """
        return f"{self._special_character}{key}{self._special_character}"