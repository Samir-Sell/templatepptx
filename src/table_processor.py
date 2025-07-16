from parent_processor import ParentProcessor
from copy import deepcopy
from pptx.table import _Cell, Table
import warnings
from text_processor import TextProcessor
from  pptx.shapes.autoshape import Shape
from typing import Collection
from typing import Union
from template_pptx_options import TemplatePptxOptions

class TableFailedToPopulate(Exception):
    """Raised when a table fails to populate due to a data or logic issue."""
    def __init__(self, message="The table could not be processed.", *, cause=None):
        super().__init__(message)
        self.__cause__ = cause 


class TableProcessor(ParentProcessor):

    def __init__(self, shape: Shape, context: dict, slide_number: int, special_character: str):
        super().__init__(shape, context, slide_number, special_character)

    
    def _remove_row(self, table: Table, row_num: int) -> None:

        """
        Description: Remove a row from a PPTX table that is found in a pptx shape object

        @input table: A table from a PPTX
        @input row_num: An integer containing the row number to be deleted. negative numbers can be used to index 
        backwords
        """
        table._tbl.remove(table._tbl.tr_lst[row_num])

    def _add_row(self, table: Table, relationship_class: str, record_count: int) -> None:

        """
        Description: Add a row to a PPTX table using a relationship specified in the contect dictionary through dot notation ex. (relaltionship_name.field_value)

        @input table: A table from a PPTX
        @input context: A dictionary containing all of the data that is fed into the template. It contains the data 
        and the magic keywords.
        @input relationship_class: The name of the relationship class
        @input record_count: An integer specifying which item in the related records context relationship value 
        should be inserted into the cells. AKA: This selects a single dictionary from the list of dictionaries.

        """
        
        # duplicating last row of the table as a new row to be added
        new_row = deepcopy(table._tbl.tr_lst[1])
        # Loop through cells in new row
        for tc in new_row.tc_lst:
            # Get cell
            cell = _Cell(tc, new_row.tc_lst)

            # Get cell text
            cell_text = cell.text
            # Process and get rel name and field in rel
            cell_value = cell_text.replace(self._special_character, "")
            cell_keys = cell_value.split(".")
            
            cell_field = cell_keys[1] # Get relationship name
            cell_field = cell_field.strip('\n') # Strip out newlines if needed

            # Replace the placeholder formatting text with text from context
            cell_text = cell_text.replace(cell_text, self._context[relationship_class][record_count][cell_field])

            # Access the first paragraph
            p = cell.text_frame.paragraphs[0]

                # Try to preserve formatting if a run exists
            if p.runs:
                original_run = p.runs[0]
                original_font = original_run.font
            else:
                original_font = None

            p.clear()
            # cell.text = cell_text        

            run = p.add_run()
            run.text = cell_text

            if original_font:
                run.font.name = original_font.name
                run.font.size = original_font.size
                run.font.bold = original_font.bold
                run.font.italic = original_font.italic
                if original_font.color and hasattr(original_font.color, 'rgb'):
                    run.font.color.rgb = original_font.color.rgb
                elif original_font.color and hasattr(original_font.color, 'theme_color'):
                    run.font.color.theme_color = original_font.color.theme_color

                    

        table._tbl.append(new_row) #Append to existing table
        
    def process_table(self, options: TemplatePptxOptions):

        '''
        Description: Process a table and replace every value or populate a table based on a relationship
        
        @input shape: A container from PPTX Python called shape which contains paragraphs and text
        @input context: A dictionary containing all of the data that is fed into the template. It contains the data 
        and the magic keywords.
        @input slide_number: The slide number index
        '''

        try:
            table_cells = self._shape.table.iter_cells()
            self._process_table_cells(table_cells)
        except Exception as e:
            if options.strict_mode:
                raise TableFailedToPopulate("Failed while processing table.", cause=e) from e
            warnings.warn(f"Table failed to be populated due to {e}")

    def _process_relationship(self, relationship_class: str, rel_class_key: str) -> Union[int, None]:
        """
        Description: Process a relationship in a table and replace text with context values 
        
        @input relationship_class: The name of the relationship class
        @input rel_class_key: The key for the relationship class in the context dictionary
        """
        if rel_class_key: 
            for row in range(len(self._context[relationship_class])):
                self._add_row(self._shape.table, relationship_class, row)
            self._remove_row(self._shape.table, 1)
            return -1
        else:
            warnings.warn(f"Relationship link for {relationship_class} does not exist.")
            raise KeyError(relationship_class)

    def _process_cell(self, cell: _Cell) -> Union[int, None]:
        """
        Description: Process a single cell in a table and replace text with context values or populate a table based on a relationship
        
        @input cell: A cell from a PPTX table
        """
        if (cell.text.find("relationship")) != -1:
            cleaned_cell = (cell.text).replace(self._special_character, "")
            relationship_class = (cleaned_cell).split(".")[0]
            rel_class_key = self._context.get(relationship_class)
            if self._process_relationship(relationship_class, rel_class_key) == -1:
                return -1
        for p in cell.text_frame.paragraphs:
            TextProcessor(self._shape, self._context, self._slide_number, self._special_character)._replace_runs(p)

    def _process_table_cells(self, table_cells: Collection[_Cell]) -> None:
        '''
        Description: Process an array of table cells
        
        @input table_cells: A collection of table cells from a PPTX table
        '''
        for cell in table_cells:
            result = self._process_cell(cell)
            if result == -1:
                return