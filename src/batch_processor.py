import os
import glob
import copy

from pptx import Presentation
from pptx.util import Inches
import tempfile

class batchTool():

    def __init__(self, pptx_dir, output_pptx):
        self._pptx_dir = pptx_dir
        self._output_pptx = output_pptx
    
    def _sort_by_number_file_names(self, in_string):

        """
        Description: Sorts a list of file names by their mumeric names. 
        Function assumes filename only contains number characters.

        @input: A file string path which contains a file with only a numeric name. Ex. 'C:\\Users\\Bob\\AppData\\Local\\Temp\\tmp0cnkqzuj\\9.pptx'
        
        """

        file_name = os.path.basename(in_string)

        if os.path.splitext(file_name)[0].isnumeric():
            return int(os.path.splitext(file_name)[0])
        else:
            return in_string

    def combine_slides(self):

        """
        Description: Combines slides from multiple PowerPoints into one PowerPoint File. This function
        does not use PowerPoint. However, it has limited functionality and is restricted to pictures,
        text and tables.

        @input pres_dir: A folder path to directory containing all of the PPTX
        @input final_output: A file string that specifies where the final combined output PowerPoint is written

        """

        # Find all slides in the temp output dir
        pres = glob.glob(os.path.join(self._pptx_dir,"*.pptx"))
        pres.sort(key=self._sort_by_number_file_names)
        #UNCOMMENT WHEN ADDING ABILITY TO ADD TEMPLATES
        #combined_presentation = Presentation(os.path.join(SCRIPT_DIR, "config_templates\\blankPres.pptx"))
        combined_presentation = Presentation()
        combined_presentation.slide_width = Inches(13.333)
        combined_presentation.slide_height = Inches(7.5)
        for presentation in pres:
            pres = Presentation(presentation)
            for slide in pres.slides:
                combined_slide = combined_presentation.slides.add_slide(combined_presentation.slide_layouts[6])
                for shape in slide.shapes:
                    if shape.shape_type == 17: # Text
                        element = copy.deepcopy(shape.element)
                        combined_slide.shapes._spTree.insert_element_before(element, 'p:extLst')
                    elif shape.shape_type == 19: # Table
                        element = copy.deepcopy(shape.element)
                        combined_slide.shapes._spTree.insert_element_before(element, 'p:extLst')
                    elif shape.shape_type == 13: # Image
                        self._replace_picture_pptx(shape, combined_slide)
                    elif shape.shape_type == 1: #  Autoshape
                        element = copy.deepcopy(shape.element)
                        combined_slide.shapes._spTree.insert_element_before(element, 'p:extLst')
                    else:
                        element = copy.deepcopy(shape.element)
                        combined_slide.shapes._spTree.insert_element_before(element, 'p:extLst')

        combined_presentation.save(self._output_pptx)

    def _replace_picture_pptx(self, shape, slide):

        """
        Description: The function to replace an image in the PowerPoint Template that does not use context
        @input shape: A container which has a image attribute associated
        @input slide: Slide object which will have the picture added to it
        """
        # Get info about template picture in order to mimic it
        img_width = shape.width
        img_height = shape.height
        img_left = shape.left
        img_top = shape.top

        blob = shape.image._blob
        with tempfile.TemporaryFile() as image:
            image.write(blob)
            slide.shapes.add_picture(image_file=image, left=img_left, top=img_top, width=img_width, height=img_height)
        