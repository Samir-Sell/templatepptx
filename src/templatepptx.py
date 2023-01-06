import warnings
from pptx import Presentation
from pptx.util import Inches
import os
import glob
import copy
import tempfile

# Custom
from text_factory import textProcessor
from table_factory import tableProcessor
from picture_factory import pictureProcessor

class templatePptx:
 
    def __init__(self, ppt, context, output_path, special_character="$"):
        self._ppt = ppt
        self._context = context
        self._output_path = output_path
        self._validation()
        self._special_character = special_character
        self.parse_template_pptx()


    def _validation(self):
        
        # Find template pptx
        try:
            self._ppt = Presentation(self._ppt)
        except Exception as e:
            raise FileNotFoundError(f"{e}")

        # Warn user if context obj is empty
        if self._context == {}:
            warnings.warn("Context file is empty")
        
        # Check if the context is a valid dictionary
        if not isinstance(self._context, dict):
            raise ValueError(f"Your context is not a valid dictionary. Please check the parameter type.")

        # Check if the output is valid
        try:
            with open(self._output_path, 'w') as out_pptx:
                pass
        except Exception as e:
            raise IOError(f"Cannot open a PPTX file at the desired output dir: {e}")


    def parse_template_pptx(self) -> Presentation:
        
        """
        Description: The parent method that parses the powerpoint into a PPTX Presentation and replaces magic words

        @input ppt: A file path to the template PPTX
        @input context: A dictionary containing all of the data that is fed into the template. It contains the data 
        and the magic keywords.

        @output ppt: A Python pptx Presentation object which contains all of the new changes
        """

        # Loop through every shape element in each slide and replace template words with values from context
        for slide in self._ppt.slides:
            slide_number = (self._ppt.slides.index(slide)) + 1
            shapes_on_slide = slide.shapes
            for shape in shapes_on_slide:
                # Process all text on the shape
                textProcessor(shape, self._context, slide_number, self._special_character).replace_text()               
                # If shape object has a table associated, process table
                # NOTE: relationship is a key word and is used to specify table relates                      
                if shape.has_table:
                    tableProcessor(shape, self._context, slide_number, self._special_character).process_table()
                # 13 is the shapetype for an image
                if shape.shape_type == 13:
                    pictureProcessor(shape, self._context, slide_number, slide).replace_picture()
                # 6 is a group shape
                if shape.shape_type == 6:
                    for sub_shape in shape.shapes:
                        textProcessor(sub_shape, self._context, slide_number, self._special_character).replace_text()
                        if sub_shape.shape_type == 13:
                            pictureProcessor(shape, self._context, slide_number, slide).replace_picture()
        self._ppt.save(self._output_path)
        return self._output_path


class batchTool():

    def __init__(self, pptx_dir, output_pptx):
        self._pptx_dir = pptx_dir
        self._output_pptx = output_pptx
    
    def _sort_by_number_file_names(self, in_string):

        """
        Description: Sorts a list of file names by their mumeric names. 
        Function assumes filename only contains number characters.

        @input: A file string path which contains a file with only a numeric name. Ex. 'C:\\Users\\Bob\\AppData\\Local\\Temp\\tmp0cnkqzuj\\9.pptx'
        
        """

        file_name = os.path.basename(in_string)

        if os.path.splitext(file_name)[0].isnumeric():
            return int(os.path.splitext(file_name)[0])
        else:
            return in_string

    def combine_slides(self):

        """
        Description: Combines slides from multiple PowerPoints into one PowerPoint File. This function
        does not use PowerPoint. However, it has limited functionality and is restricted to pictures,
        text and tables.

        @input pres_dir: A folder path to directory containing all of the PPTX
        @input final_output: A file string that specifies where the final combined output PowerPoint is written

        """

        # Find all slides in the temp output dir
        pres = glob.glob(os.path.join(self._pptx_dir,"*.pptx"))
        pres.sort(key=self._sort_by_number_file_names)
        #UNCOMMENT WHEN ADDING ABILITY TO ADD TEMPLATES
        #combined_presentation = Presentation(os.path.join(SCRIPT_DIR, "config_templates\\blankPres.pptx"))
        combined_presentation = Presentation()
        combined_presentation.slide_width = Inches(13.333)
        combined_presentation.slide_height = Inches(7.5)
        for presentation in pres:
            pres = Presentation(presentation)
            for slide in pres.slides:
                combined_slide = combined_presentation.slides.add_slide(combined_presentation.slide_layouts[6])
                for shape in slide.shapes:
                    if shape.shape_type == 17: # Text
                        element = copy.deepcopy(shape.element)
                        combined_slide.shapes._spTree.insert_element_before(element, 'p:extLst')
                    elif shape.shape_type == 19: # Table
                        element = copy.deepcopy(shape.element)
                        combined_slide.shapes._spTree.insert_element_before(element, 'p:extLst')
                    elif shape.shape_type == 13: # Image
                        self._replace_picture_pptx(shape, combined_slide)
                    elif shape.shape_type == 1: #  Autoshape
                        element = copy.deepcopy(shape.element)
                        combined_slide.shapes._spTree.insert_element_before(element, 'p:extLst')
                    else:
                        element = copy.deepcopy(shape.element)
                        combined_slide.shapes._spTree.insert_element_before(element, 'p:extLst')

        combined_presentation.save(self._output_pptx)

    def _replace_picture_pptx(self, shape, slide):

        """
        Description: The function to replace an image in the PowerPoint Template that does not use context
        @input shape: A container which has a image attribute associated
        @input slide: Slide object which will have the picture added to it
        """
        # Get info about template picture in order to mimic it
        img_width = shape.width
        img_height = shape.height
        img_left = shape.left
        img_top = shape.top

        blob = shape.image._blob
        with tempfile.TemporaryFile() as image:
            image.write(blob)
            slide.shapes.add_picture(image_file=image, left=img_left, top=img_top, width=img_width, height=img_height)
        