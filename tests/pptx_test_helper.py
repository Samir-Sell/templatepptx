from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from PIL import Image, ImageDraw

def _generate_image(path: str, text: str, size=(300, 200), color=(100, 100, 200)):
    Path(path).parent.mkdir(parents=True, exist_ok=True)

    image = Image.new("RGB", size, color=color)
    draw = ImageDraw.Draw(image)
    draw.text((10, 10), text, fill=(255, 255, 255))
    image.save(path)


def create_photo_presentation(output_pptx="tests/templates/photo_test.pptx"):
    # Paths
    placeholder_img = "tests/assets/placeholder.png"
    real_img = "tests/assets/photo1.png"

    # Generate images
    _generate_image(placeholder_img, "Placeholder", color=(150, 150, 150))
    _generate_image(real_img, "Actual Photo", color=(200, 50, 50))

    # Create PowerPoint template
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide

    left = Inches(2)
    top = Inches(1.5)
    height = Inches(3)

    picture = slide.shapes.add_picture(placeholder_img, left, top, height=height)
    picture.name = "Photo Placeholder"

    Path(output_pptx).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)

def create_text_presentation(output_path="tests/templates/textbox_test.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 

    example_texts = [
        {"text": "Here is $exampleone$ in bold red.", "font_size": 32, "color": RGBColor(255, 0, 0), "bold": True},
        {"text": "Now $exampletwo$ appears in italic blue.", "font_size": 28, "color": RGBColor(0, 0, 255), "italic": True},
        {"text": "$examplethree$ is big and green.", "font_size": 44, "color": RGBColor(0, 128, 0)},
        {"text": "Another mention: $examplefour$ with underline.", "font_size": 26, "color": RGBColor(128, 0, 128), "underline": True},
        {"text": "Crazy font mix $examplefive$ in orange.", "font_size": 20, "color": RGBColor(255, 165, 0)},
        {"text": "$examplesix$ in small gray text.", "font_size": 16, "color": RGBColor(100, 100, 100)},
    ]

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, style in enumerate(example_texts):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        run = p.add_run()
        run.text = style["text"]
        font = run.font
        font.size = Pt(style["font_size"])
        font.color.rgb = style["color"]
        font.bold = style.get("bold", False)
        font.italic = style.get("italic", False)
        font.underline = style.get("underline", False)
        p.alignment = PP_ALIGN.LEFT

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

def create_table_presentation(output_path="tests/templates/table_test.pptx"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

    rows = 2  # Header + one magic row
    cols = 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Optional: Set column widths for readability
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)

    # Header
    table.cell(0, 0).text = "ID"
    table.cell(0, 1).text = "First Name"
    table.cell(0, 2).text = "Last Name"

    # Magic word row (correct format)
    table.cell(1, 0).text = "$relationship_people.id$"
    table.cell(1, 1).text = "$relationship_people.first_name$"
    table.cell(1, 2).text = "$relationship_people.last_name$"

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    