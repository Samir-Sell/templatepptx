import unittest
from pathlib import Path
from pptx import Presentation
from pathlib import Path
from pptx_test_helper import create_text_presentation, create_table_presentation, create_photo_presentation

import sys
sys.path.append("src")
from templatepptx import TemplatePptx

# Define template and output folders
TEMPLATE_DIR = Path("tests/templates")
OUTPUT_DIR = Path("tests/outputs")
ASSETS_DIR = Path("tests/assets")



class TemplateTestHelper:
    """Reusable helper for PowerPoint magic word testing."""

    def __init__(self, template_file, output_file, context, special_character="$"):
        self.template_path = TEMPLATE_DIR / template_file
        self.output_path = OUTPUT_DIR / output_file
        self.context = context
        self.special_character = special_character

    def run_template_engine(self):
        ppt = TemplatePptx(
            str(self.template_path),
            self.context,
            str(self.output_path),
            self.special_character
        )
        ppt.options.strict_mode = True
        ppt.parse_template_pptx()

    def extract_text(self):
        prs = Presentation(str(self.output_path))
        found_text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                # Text box, title, etc.
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            found_text.append(run.text)

                # Table
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            found_text.append(cell.text)

        return " ".join(found_text)
    
    def extract_image_alt_texts(self):
        prs = Presentation(str(self.output_path))
        alt_texts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    text = shape.element.xpath("//p:pic/p:nvPicPr/p:cNvPr")[0].attrib["descr"]
                    alt_texts.append(text)
        return alt_texts


    def assert_replacements(self, test_case):
        text = self.extract_text()
        for key, value in self.context.items():
            test_case.assertIn(value, text)
            test_case.assertNotIn(f"${key}$", text)


class TestTemplatePptx(unittest.TestCase):

    @classmethod
    def setUpClass(cls):
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        create_text_presentation()
        create_table_presentation()
        create_photo_presentation()

    @classmethod
    def tearDownClass(cls):
        for file in OUTPUT_DIR.glob("*.pptx"):
            file.unlink(missing_ok=True)

        for file in TEMPLATE_DIR.glob("*.pptx"):
            file.unlink(missing_ok=True)

        for file in ASSETS_DIR.glob("*.png"):
            file.unlink(missing_ok=True)
    

    def test_textbox_replacement(self):

        context = {
        "exampleone": "filled_one",
        "exampletwo": "filled_two",
        "examplethree": "filled_three",
        "examplefour": "filled_four",
        "examplefive": "filled_five",
        "examplesix": "filled_six"
        }

        helper = TemplateTestHelper(
            template_file="textbox_test.pptx",
            output_file="textbox_output.pptx",
            context=context
        )
        helper.run_template_engine()
        helper.assert_replacements(self)

    def test_table_replacement(self):
        context = {
            "relationship_people": [
                {"id": "1", "first_name": "Alice", "last_name": "Anderson"},
                {"id": "2", "first_name": "Bob", "last_name": "Brown"},
                {"id": "3", "first_name": "Charlie", "last_name": "Clark"}
            ]
        }

        helper = TemplateTestHelper(
            template_file="table_test.pptx",
            output_file="table_output.pptx",
            context=context
        )
        helper.run_template_engine()

        text = helper.extract_text()

        # Check all values in each row were filled in
        for person in context["relationship_people"]:
            for value in person.values():
                assert value in text, f"Expected value '{value}' not found in slide text"

        # Ensure placeholders are gone
        assert "$relationship_people.id$" not in text
        assert "$relationship_people.first_name$" not in text
        assert "$relationship_people.last_name$" not in text

    def test_photo_replacement(self):
        context = {
            "placeholder.png": "tests/assets/photo1.png"
        }

        helper = TemplateTestHelper(
            template_file="photo_test.pptx",
            output_file="photo_output.pptx",
            context=context
        )
        helper.run_template_engine()

        alt_texts = helper.extract_image_alt_texts()
        print(alt_texts)
        for alt in alt_texts:
            self.assertNotIn("placeholder.png", alt, "Placeholder alt text was not replaced in image")
        

if __name__ == "__main__":
    unittest.main()
