import unittest
import pptx
import os
import shutil
import json
import picture_factory, table_factory, text_factory, templatepptx


class test_text_factory(unittest.TestCase):

    # Classmethod run once after all tests are completed
    @classmethod
    def tearDownClass(cls):
        pass

    # Run before every single unit test
    def setUp(self):

        # Loads PowerPoint into a PPTX Obj
        SCRIPT_PATH = os.path.dirname(os.path.realpath(__file__))
        self.pptx_unit_test_path = os.path.join(SCRIPT_PATH, r"testInputs/UnitTesting.pptx")
        self.ppt = pptx.Presentation(self.pptx_unit_test_path)


    def test_formatting_check(self):

        # Tests the formatting_check to ensure valid run detection 
        # Uses slide 1 and 2

        tft = text_factory.textProcessor(self.ppt.slides[0].shapes[0],{},1,"$")

        # Returns run False due to None and due to bold making it return false
        paragraphs = self.ppt.slides[0].shapes[0].text_frame.paragraphs
        last_run = None
        for p in paragraphs:
            for run in p.runs:
                self.assertFalse(tft._formatting_check(run, last_run))

        # Expected to return True due to no formatting differenes
        paragraphs = self.ppt.slides[1].shapes[0].text_frame.paragraphs
        last_run = self.ppt.slides[1].shapes[0].text_frame.paragraphs[0].runs[0]
        for p in paragraphs:
            for run in p.runs:
                self.assertTrue(tft._formatting_check(run, last_run))


    def test_replace_runs(self):

        # Uses slides 3 and 4
        # Test if varying run values are replaced correctly 

        # Test to ensure correct replacement
        tft = text_factory.textProcessor(self.ppt.slides[2].shapes[0],{"test":"replaced"},3,"$")
        paragraph = self.ppt.slides[2].shapes[0].text_frame.paragraphs[0]
        tft._replace_runs(paragraph)
        assert self.ppt.slides[2].shapes[0].text == "replaced"

        # Test to ensure correct replacement even when run contains multiple formatting
        tft = text_factory.textProcessor(self.ppt.slides[3].shapes[0],{"test":"replaced"},4,"$")
        paragraph = self.ppt.slides[3].shapes[0].text_frame.paragraphs[0]
        tft._replace_runs(paragraph)
        assert self.ppt.slides[2].shapes[0].text == "replaced"

    def test_add_special(self):
        
        # No slides used. tft inititalized only for test purpose
        tft = text_factory.textProcessor(self.ppt.slides[3].shapes[0],{"test":"replaced"},4,"$")
        self.assertEqual(tft._add_special(key="test"), "$test$")



class test_table_factory(unittest.TestCase):

    # Classmethod run once after all tests are completed
    @classmethod
    def tearDownClass(cls):
        pass

    # Run before every single unit test
    def setUp(self):

        # Loads PowerPoint into a PPTX Obj
        SCRIPT_PATH = os.path.dirname(os.path.realpath(__file__))
        self.pptx_unit_test_path = os.path.join(SCRIPT_PATH, r"testInputs/UnitTesting.pptx")
        self.ppt = pptx.Presentation(self.pptx_unit_test_path)

    def test_remove_row(self):

        # Uses slide 5
        # Removes a row from the table on slide 5 and confirms new number of rows
        ttf = table_factory.tableProcessor(self.ppt.slides[4].shapes[0], {}, 5, "$")

        table = self.ppt.slides[4].shapes[0].table
        ttf._remove_row(table, 1)
        assert len(self.ppt.slides[4].shapes[0].table.rows) == 2


    def test_add_row(self):
        
        # Uses slide 6
        # Adds rows from context to table on slide 6 and confirms relationships are working correctly

        table = self.ppt.slides[5].shapes[0].table
        relationship_class = "relationship_test"
        record_count = 1

        ttf = table_factory.tableProcessor(self.ppt.slides[5].shapes[0], 
                                           {"relationship_test":[{"test1": "a", "test2": "c", "test3": "c"},
                                                                 {"test1": "b", "test2": "c", "test3": "c"}]},
                                           6, "$")

        ttf._add_row(table, relationship_class, record_count)
        assert len(table.rows) == 3
        ttf._add_row(table, relationship_class, record_count)
        assert len(table.rows) == 4


    def test_process_table(self):

        # Uses slides 11 and 12
        
        # Test a table with relatonship values in slide 11. Add values and then check if values exist in the table
        shape = self.ppt.slides[10].shapes[0]
        context = {"relationship_test":[{"test1": "a", "test2": "b", "test3": "c"},{"test1": "d", "test2": "e", "test3": "f"}]}
        slide_number = 6

        ttf = table_factory.tableProcessor(self.ppt.slides[10].shapes[0], context, slide_number, "$")
        ttf.process_table()
        rows = self.ppt.slides[10].shapes[0].table.rows
        self.assertEqual(len(rows), 3)
        expected_values = [ "test0", "test1", "test2", "a", "b", "c", "d", "e", "f" ]
        for cell in shape.table.iter_cells():
            assert cell.text in expected_values

        # Test table with no relationship values in slide 12 and confirm they are replacec correctly
        shape = self.ppt.slides[11].shapes[0]
        context = {"testrow0": "I am first column", "testrow1": "I am second column", "testrow2": "I am third column"}
        slide_number = 12
        ttf = table_factory.tableProcessor(self.ppt.slides[11].shapes[0], context, slide_number, "$")

        ttf.process_table()
        rows = self.ppt.slides[11].shapes[0].table.rows
        self.assertEqual(len(rows), 2)
        expected_values = [ "test0", "test1", "test2", "I am first column", "I am second column", "I am third column"]
        for cell in shape.table.iter_cells():
            assert cell.text in expected_values


class test_picture_factory(unittest.TestCase):

    # Classmethod run once after all tests are completed
    @classmethod
    def tearDownClass(cls):
        pass

    # Run before every single unit test
    def setUp(self):

        # Loads PowerPoint into a PPTX Obj
        SCRIPT_PATH = os.path.dirname(os.path.realpath(__file__))
        self.pptx_unit_test_path = os.path.join(SCRIPT_PATH, r"testInputs/UnitTesting.pptx")
        self.ppt = pptx.Presentation(self.pptx_unit_test_path)

    def test_get_alt_text(self):

        # Uses slides 9 and 10
        slide_number = 9
        shape = self.ppt.slides[8].shapes[0]
        tpf = picture_factory.pictureProcessor(shape, {}, slide_number, self.ppt.slides[8], "$")
        
        # Test to find correct alt text from image
        self.assertEqual(tpf._get_alt_text(), "unit_testing_example")
        
        # Test to find empty alt text from images
        slide_number = 10
        shape = self.ppt.slides[9].shapes[0]
        tpf = picture_factory.pictureProcessor(shape, {}, slide_number, self.ppt.slides[9], "$")

        self.assertEquals(tpf._get_alt_text(), None)
    

if __name__ == '__main__':
    unittest.main()

